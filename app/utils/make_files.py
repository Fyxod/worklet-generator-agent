gap =0.3
import asyncio
import os
import re

from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas
from reportlab.platypus import Frame, Paragraph

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.utils.reference_functions.reference_sort import index_sort

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE_DIR)

ppt_path = os.path.join(PROJECT_ROOT, "resources/generated_worklets/ppt")
pdf_path = os.path.join(PROJECT_ROOT, "resources/generated_worklets/pdf")

os.makedirs(pdf_path, exist_ok=True)
os.makedirs(ppt_path, exist_ok=True)


CUSTOM_PAGE_SIZE = (750,900)  # Width x Height in points (1 point = 1/72 inch)  used by pdf

async def generatePdf(json, model, index):

    json_data = await pre_processing(json, index, model) # json after sorting
    
    safe_title = sanitize_filename(json.get("Title", "untitled"))
    
    filename_pdf = os.path.join(pdf_path, f"{safe_title}.pdf")
    filename_ppt = os.path.join(ppt_path, f"{safe_title}.pptx")

    loop = asyncio.get_running_loop()
    await loop.run_in_executor(None, create_pdf, filename_pdf, json_data)
    await loop.run_in_executor(None, create_ppt, filename_ppt, json_data)
    
async def pre_processing(json_data, index, model="gemma3:27b"):
    json_data = await index_sort(json_data, model)
    return json_data

def sanitize_filename(filename):
    return re.sub(r'[\/:*?"<>|]', '_', filename)

# Function to handle the actual blocking PDF generation using ReportLab
def create_pdf(filename, json_data):
    pdf = canvas.Canvas(filename, pagesize=CUSTOM_PAGE_SIZE)
    width, height = CUSTOM_PAGE_SIZE

    # Setup the styles for the PDF content
    styles = getSampleStyleSheet()
    header_style = ParagraphStyle('header_style', parent=styles['Heading1'], fontSize=20, textColor=colors.darkblue)
    normal_style = ParagraphStyle('normal_style', parent=styles['BodyText'], fontSize=12, leading=15)
    bullet_style = ParagraphStyle('bullet_style', parent=styles['BodyText'], fontSize=12, leftIndent=20, bulletIndent=10)

    y = height - 50
    frame = Frame(40, 40, width - 80, height - 100, showBoundary=0)

    elements = []

    # Add Title and other fields from json_data
    if 'Title' in json_data:
        elements.append(Paragraph(f"<b>Title:</b> {json_data['Title']}", header_style))
    if 'Problem Statement' in json_data:
        elements.append(Paragraph(f"<b>Problem Statement:</b> {json_data['Problem Statement']}", normal_style))
    if 'Description' in json_data:
        elements.append(Paragraph(f"<b>Description:</b> {json_data['Description']}", normal_style))
    if 'Challenge / Use Case' in json_data:
        elements.append(Paragraph(f"<b>Challenge / Use Case:</b> {json_data['Challenge / Use Case']}", normal_style))
    if 'Deliverables' in json_data:
        elements.append(Paragraph(f"<b>Deliverables:</b> {json_data['Deliverables']}", normal_style))

    if 'KPIs' in json_data and isinstance(json_data['KPIs'], list):
        elements.append(Paragraph("<b>KPIs:</b>", normal_style))
        for kpi in json_data['KPIs']:
            elements.append(Paragraph(f"• {kpi}", bullet_style))

    if 'Prerequisites' in json_data and isinstance(json_data['Prerequisites'], list):
        elements.append(Paragraph("<b>Prerequisites:</b>", normal_style))
        for prereq in json_data['Prerequisites']:
            elements.append(Paragraph(f"• {prereq}", bullet_style))

    if 'Infrastructure Requirements' in json_data:
        elements.append(Paragraph(f"<b>Infrastructure Requirements:</b> {json_data['Infrastructure Requirements']}", normal_style))
    if 'Tentative Tech Stack' in json_data:
        elements.append(Paragraph(f"<b>Tentative Tech Stack:</b> {json_data['Tentative Tech Stack']}", normal_style))

    if 'Milestones (6 months)' in json_data and isinstance(json_data['Milestones (6 months)'], dict):
        elements.append(Paragraph("<b>Milestones (6 months):</b>", normal_style))
        milestones = json_data['Milestones (6 months)']
        if 'M2' in milestones:
            elements.append(Paragraph(f"• M2: {milestones['M2']}", bullet_style))
        if 'M4' in milestones:
            elements.append(Paragraph(f"• M4: {milestones['M4']}", bullet_style))
        if 'M6' in milestones:
            elements.append(Paragraph(f"• M6: {milestones['M6']}", bullet_style))

    if 'Reference Work' in json_data and isinstance(json_data['Reference Work'], list):
        elements.append(Paragraph("<b>Reference Work:</b>", normal_style))
        for idx, ref in enumerate(json_data['Reference Work']):
            if isinstance(ref, dict) and 'Title' in ref and 'Link' in ref:
                link_paragraph = f'<a href="{ref["Link"]}">{ref["Title"]}</a>'
                elements.append(Paragraph(link_paragraph, bullet_style))

    frame.addFromList(elements, pdf)
    pdf.save()

    print(f"PDF generated: {filename}")


def create_ppt(output_filename, json_data):
    data = [json_data]
    prs = Presentation()
    prs.slide_width = Pt(750)
    prs.slide_height = Pt(1100)
    slide_layout = prs.slide_layouts[6]
    for entry in data:
        slide = prs.slides.add_slide(slide_layout)
        top = 0.5

        top = add_textbox_Title(slide, "Title", entry.get("Title", ""), top)+0.1
        top = add_textbox(slide, "Problem Statement", entry.get("Problem Statement", ""), top)
        top = add_textbox(slide, "Description", entry.get("Description", ""), top)
        top = add_textbox(slide, "Challenge / Use Case", entry.get("Challenge / Use Case", ""), top)
        top = add_textbox(slide, "Deliverables", entry.get("Deliverables", ""), top)

        kpis = entry.get("KPIs", [])
        if kpis:
            kpi_text = "\n".join([f"• {k}" for k in kpis])
            top = add_textbox(slide, "KPIs", kpi_text, top)

        preq = entry.get("Prerequisites", [])
        if preq:
            preq_text = "\n".join([f"• {p}" for p in preq])
            top = add_textbox(slide, "Prerequisites", preq_text, top)

        top = add_textbox(slide, "Infrastructure Requirements", entry.get("Infrastructure Requirements", ""), top)
        top = add_textbox(slide, "Tentative Tech Stack", entry.get("Tentative Tech Stack", ""), top)

        milestones = entry.get("Milestones (6 months)", {})
        if milestones:
            milestone_text = "\n".join([f"{k}: {v}" for k, v in milestones.items()])
            top = add_textbox(slide, "Milestones (6 months)", milestone_text, top)

        references = entry.get("Reference Work", [])
        if references:
            left = Inches(0.5)
            top_inch = Inches(top)
            width = Inches(9.5)
            height = Inches(len(references) * 0.4 + 0.2)
            textbox = slide.shapes.add_textbox(left, top_inch, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()

            title_para = tf.paragraphs[0]
            title_run = title_para.add_run()
            title_run.font.size = Pt(16)
            title_run.font.bold = True
            title_run.font.name = 'Calibri'
            title_run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
            title_run.text = "Reference Work:"

            for ref in references:
                p = tf.add_paragraph()
                p.level = 1
                run = p.add_run()
                run.text = f"• {ref['Title']}"
                run.font.size = Pt(15)
                run.font.name = 'Calibri'
                run.font.color.rgb = RGBColor(0, 102, 204)
                run.hyperlink.address = ref['Link']

            top += height.inches + gap
    prs.save(output_filename)

def estimate_height_wrapped_content(text, chars_per_line=100, line_height_pt=18):
    lines = 0
    for para in text.split('\n'):
        para = para.strip()
        if not para:
            continue
        lines += max(1, int(len(para) / chars_per_line) + 1)
    return Pt(lines * line_height_pt).inches

def estimate_height_wrapped_Title(text, chars_per_line=65, line_height_pt=20):
    lines = 0
    for para in text.split('\n'):
        para = para.strip()
        if not para:
            continue
        lines += max(1, int(len(para) / chars_per_line) + 1)
    return Pt(lines * line_height_pt).inches

def add_textbox(slide, title, content, top_inch):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.5)
    height = estimate_height_wrapped_content(content)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run_title = p.add_run()
    run_title.text = f"{title}:\n"
    run_title.font.size = Pt(16)
    run_title.font.bold = True
    run_title.font.name = 'Calibri'
    run_title.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    run_content = p.add_run()
    run_content.text = content
    run_content.font.size = Pt(15)
    run_content.font.name = 'Calibri'

    return top_inch + height + gap

def add_textbox_Title(slide, title, content, top_inch):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.5)
    height = estimate_height_wrapped_Title(content)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run_title = p.add_run()
    run_title.text = f"{title}: "
    run_title.font.size = Pt(20)
    run_title.font.name = 'Calibri'
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    run_content = p.add_run()
    run_content.text = content
    run_content.font.size = Pt(20)
    run_content.font.name = 'Calibri'
    run_content.font.bold = True
    run_content.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    return top_inch +height+ gap

