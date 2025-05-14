import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


sample_json = {
        "Title": "Federated Learning for Personalized Healthcare in Rural India",
        "Problem Statement": "Develop a federated learning framework leveraging edge devices to train a personalized disease prediction model for rural Indian patients, addressing data privacy and limited internet connectivity.",
        "Description": "Healthcare access in rural India is limited by data silos aHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. ThisHealthcare access in rural India is limited by data silos and privacy concerns. Thisnd privacy concerns. This project aims to create a distributed AI model using federated learning, enabling healthcare providers to learn from patient data without direct data sharing, improving diagnostic accuracy and treatment planning.",
        "Challenge / Use Case": "Limited access to specialized healthcare and challenges in data sharing due to privacy regulations in rural Indian clinics.",
        "Deliverables": "A federated learning framework, a trained disease prediction model, a secure data aggregation mechanism, and a detailed performance report.",
        "KPIs": [
            "Accuracy ≥ 10% improvement over centralized model",
            "Communication Rounds ≤ 15",
            "Data Privacy: Differential Privacy budget (ε) ≤ 1.0",
            "Model Convergence: Achieve convergence within 30 epochs"
        ],
        "Prerequisites": [
            "Understanding of federated learning principles",
            "Familiarity with machine learning algorithms (e.g., logistic regression, SVM)",
            "Knowledge of privacy-preserving techniques (e.g., differential privacy)",
            "Experience with Python and TensorFlow/PyTorch",
            "Basic networking concepts",
            "Familiarity with Indian healthcare data formats (if available)"
        ],
        "Infrastructure Requirements": "Cloud credits (Google Cloud/AWS) for model training and edge device simulation; 2-4 GPUs recommended for accelerated training.",
        "Tentative Tech Stack": "Python, TensorFlow/PyTorch, Federated Learning Framework (Flower/FedML), Differential Privacy library (Opacus), Scikit-learn",
        "Milestones (6 months)": {
            "M2": "Implement basic federated learning setup with synthetic data",
            "M4": "Integrate with a simulated rural healthcare dataset and implement privacy budget",
            "M6": "Develop and deploy the federated learning model with a final performance report"
        },
        "Reference Work": [
            {
                "Title": "Case Studies in Federated Learning for Healthcare",
                "Link": "https://www.igi-global.com/viewtitle.aspx?titleid=346276"
            },
            {
                "Title": "Federated Learning: Advancing Healthcare through Collaborative Artificial Intelligence",
                "Link": "https://journals.lww.com/ijcn/fulltext/2024/01000/federated_learning__advancing_healthcare_through.15.aspx"
            },
            {
                "Title": "Federated Learning for Diabetic Retinopathy Diagnosis: Enhancing Accuracy and Generalizability in Under-Resourced Regions",
                "Link": "https://arxiv.org/abs/2411.00869"
            },
            {
                "Title": "Federated Learning: The much-needed intervention in Healthcare Informatics.",
                "Link": "https://search.ebscohost.com/login.aspx?direct=true&profile=ehost&scope=site&authtype=crawler&jrnl=26767104&AN=184661859"
            },
            {
                "Title": "Federated learning framework for consumer IoMT-edge resource recommendation under telemedicine services",
                "Link": "https://ieeexplore.ieee.org/abstract/document/10793078/"
            },
            {
                "Title": "Federated learning for the internet-of-medical-things: A survey",
                "Link": "https://www.mdpi.com/2227-7390/11/1/151"
            }
        ]
    }
# Wrap in list for consistent processing
data = [sample_json]

prs = Presentation()
prs.slide_width = Pt(750)
prs.slide_height = Pt(1100)
slide_layout = prs.slide_layouts[6]  


def estimate_height_wrapped(text, chars_per_line=90, line_height_pt=18):
    lines = 0
    for para in text.split('\n'):
        para = para.strip()
        if not para:
            continue
        # Estimate how many wrapped lines this paragraph takes
        lines += max(1, int(len(para) / chars_per_line) + 1)
    return Pt(lines * line_height_pt).inches




def add_textbox(slide, title, content, top_inch):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.5)
    height = estimate_height_wrapped(content) + 0.5
    top += height
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    p.clear()

    # Title run
    run_title = p.add_run()
    run_title.text = f"{title}:\n "
    run_title.font.size = Pt(16)  
    run_title.font.bold = True
    run_title.font.name = 'Calibri'
    run_title.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)  

    # Content run
    run_content = p.add_run()
    run_content.text = content
    run_content.font.size = Pt(15)
    run_content.font.name = 'Calibri'

def add_textbox_Title(slide, title, content, top_inch):
    left = Inches(0.5)
    top = Inches(top_inch)
    width = Inches(9.5)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    p.clear()

    # Title run
    run_title = p.add_run()
    run_title.text = f"{title}: "
    run_title.font.size = Pt(20)  
    run_title.font.name = 'Calibri'
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)  

    # Content run
    run_content = p.add_run()
    run_content.text = content
    run_content.font.size = Pt(20)
    run_content.font.name = 'Calibri'
    run_content.font.bold = True
    run_content.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)  


for entry in data:
    slide = prs.slides.add_slide(slide_layout)
    top = 0.5

    add_textbox_Title(slide, "Title", entry.get("Title", ""), top); top += .5
    add_textbox(slide, "Problem Statement", entry.get("Problem Statement", ""), top); top += 1
    add_textbox(slide, "Description", entry.get("Description", ""), top); top += 1
    add_textbox(slide, "Challenge / Use Case", entry.get("Challenge / Use Case", ""), top); top += 1
    add_textbox(slide, "Deliverables", entry.get("Deliverables", ""), top); top += 1

    # KPIs
    kpis = entry.get("KPIs", [])
    if kpis:
        kpi_text = "\n".join([f"• {k}" for k in kpis])
        add_textbox(slide, "KPIs", kpi_text, top); top += 1.1

    # Prerequisites
    preq = entry.get("Prerequisites", [])
    if preq:
        preq_text = "\n".join([f"• {p}" for p in preq])
        add_textbox(slide, "Prerequisites", preq_text, top); top += 1.0

    top +=1
    add_textbox(slide, "Infrastructure Requirements", entry.get("Infrastructure Requirements", ""), top); top += 1
    add_textbox(slide, "Tentative Tech Stack", entry.get("Tentative Tech Stack", ""), top); top += 1

    milestones = entry.get("Milestones (6 months)", {})
    if milestones:
        milestone_text = "\n".join([f"{k}: {v}" for k, v in milestones.items()])
        add_textbox(slide, "Milestones (6 months)", milestone_text, top); top += 1

    # Reference Work
    references = entry.get("Reference Work", [])
    if references:
        ref_text = "\n".join([f"• {ref['Title']} ({ref['Link']})" for ref in references])
        add_textbox(slide, "Reference Work", ref_text, top); top += 1.0

prs.save("11.pptx")